import os
from pptx import Presentation
from pptx.util import Inches

def create_presentation(base_dir, output_file):
    prs = Presentation()
    
    # Optional: set presentation slide size (e.g., 16:9)
    # Default is 4:3, for 16:9 we can set:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # We want a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    
    # Iterate through all subdirectories
    # We sort them to try and keep some logical order
    subdirs = sorted([d for d in os.listdir(base_dir) if os.path.isdir(os.path.join(base_dir, d))])
    
    for subdir in subdirs:
        dir_path = os.path.join(base_dir, subdir)
        img_path = os.path.join(dir_path, 'screen.png')
        
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # We want to fit the image on the slide.
            # We'll calculate the best fit, but for simplicity, we can add it and set height/width
            try:
                # Add image, leaving it to scale proportionally if we just specify height or width
                # Since slides are 16:9, let's make it fill height or width
                pic = slide.shapes.add_picture(img_path, 0, 0, height=prs.slide_height)
                
                # Center the picture horizontally if it doesn't fill the width
                if pic.width < prs.slide_width:
                    pic.left = int((prs.slide_width - pic.width) / 2)
                    
                print(f"Added {subdir}/screen.png")
            except Exception as e:
                print(f"Error adding {img_path}: {e}")
                
    prs.save(output_file)
    print(f"Successfully saved presentation to {output_file}")

if __name__ == '__main__':
    base_directory = r'c:\Users\Lenovo\OneDrive\Desktop\ILKOMERZ\S6\ADS\Projek\desain web\stitch_ipb_food_hub_platform'
    output_filename = r'c:\Users\Lenovo\OneDrive\Desktop\ILKOMERZ\S6\ADS\Projek\desain web\Hifi_Designs_IPB_Food_Hub.pptx'
    create_presentation(base_directory, output_filename)
